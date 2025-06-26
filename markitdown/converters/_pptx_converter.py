import base64
import html
import re
from typing import Union

import pptx
from my_langchain_community.chat_models import ChatOCIGenAI
from langchain_core.messages import HumanMessage

from ._base import DocumentConverterResult, DocumentConverter
from ._html_converter import HtmlConverter


class PptxConverter(HtmlConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """

    def __init__(
            self, priority: float = DocumentConverter.PRIORITY_SPECIFIC_FILE_FORMAT
    ):
        super().__init__(priority=priority)

    def _get_llm_description(
            self, llm_client, llm_model, image_blob, content_type, prompt=None
    ):
        print("in get_llm_description()...")
        if prompt is None or prompt.strip() == "":
            prompt = "Write a detailed alt text for this image with less than 50 words."

        image_base64 = base64.b64encode(image_blob).decode("utf-8")
        data_uri = f"data:{content_type};base64,{image_base64}"
        if isinstance(llm_client, ChatOCIGenAI):
            print("llm_client is instance of ChatOCIGenAI")
            human_message = HumanMessage(content=[
                {"type": "text", "text": prompt},
                {
                    "type": "image_url",
                    "image_url": {"url": data_uri},
                },
            ], )
            messages = [human_message]
            response = llm_client.invoke(messages)
            print(f"{response=}")
            return response.content
        else:
            print("llm_client is not instance of ChatOCIGenAI")
            messages = [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": data_uri,
                            },
                        },
                        {"type": "text", "text": prompt},
                    ],
                }
            ]

            response = llm_client.chat.completions.create(
                model=llm_model, messages=messages
            )
            print(f"{response=}")
            return response.choices[0].message.content

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a PPTX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pptx":
            return None

        md_content = ""

        presentation = pptx.Presentation(local_path)
        slide_num = 0
        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title
            for shape in slide.shapes:
                # Pictures
                if self._is_picture(shape):
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069

                    llm_description = None
                    alt_text = None

                    llm_client = kwargs.get("llm_client")
                    llm_model = kwargs.get("llm_model")
                    print(f"{llm_client=}")
                    print(f"{llm_model=}")
                    if llm_client is not None and llm_model is not None:
                        try:
                            llm_description = self._get_llm_description(
                                llm_client,
                                llm_model,
                                shape.image.blob,
                                shape.image.content_type,
                                prompt=kwargs.get("llm_prompt"),
                            ).strip()
                            + "\n"
                        except Exception:
                            # Unable to describe with LLM
                            pass

                    if not llm_description:
                        try:
                            alt_text = shape._element._nvXxPr.cNvPr.attrib.get(
                                "descr", ""
                            )
                        except Exception:
                            # Unable to get alt text
                            pass

                    # A placeholder name
                    filename = re.sub(r"\W", "", shape.name) + ".jpg"
                    md_content += (
                            # "\n!["
                            "\n["
                            + (llm_description or alt_text or shape.name)
                            + "]("
                            + filename
                            + ")\n"
                    )

                # Tables
                if self._is_table(shape):
                    html_table = "<html><body><table>"
                    first_row = True
                    for row in shape.table.rows:
                        html_table += "<tr>"
                        for cell in row.cells:
                            if first_row:
                                html_table += "<th>" + html.escape(cell.text) + "</th>"
                            else:
                                html_table += "<td>" + html.escape(cell.text) + "</td>"
                        html_table += "</tr>"
                        first_row = False
                    html_table += "</table></body></html>"
                    md_content += (
                            "\n" + self._convert(html_table).text_content.strip() + "\n"
                    )

                # Charts
                if shape.has_chart:
                    md_content += self._convert_chart_to_markdown(shape.chart)

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    def _convert_chart_to_markdown(self, chart):
        md = "\n\n### Chart"
        if chart.has_title:
            md += f": {chart.chart_title.text_frame.text}"
        md += "\n\n"
        data = []
        category_names = [c.label for c in chart.plots[0].categories]
        series_names = [s.name for s in chart.series]
        data.append(["Category"] + series_names)

        for idx, category in enumerate(category_names):
            row = [category]
            for series in chart.series:
                row.append(series.values[idx])
            data.append(row)

        markdown_table = []
        for row in data:
            markdown_table.append("| " + " | ".join(map(str, row)) + " |")
        header = markdown_table[0]
        separator = "|" + "|".join(["---"] * len(data[0])) + "|"
        return md + "\n".join([header, separator] + markdown_table[1:])
